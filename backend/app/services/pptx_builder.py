"""
Build a .pptx pitch deck from structured slide content (title + bullets per slide).
Used by Asset Forge after Nova generates the slide outline.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Any, List, Tuple

SlideResult = Tuple[List[dict] | None, str | None]  # (slides, golden_thread)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

try:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _CHARTS_AVAILABLE = True
except ImportError:
    _CHARTS_AVAILABLE = False
    CategoryChartData = None  # type: ignore
    XL_CHART_TYPE = None  # type: ignore


def _add_market_chart_placeholder(slide, left, top, width, height) -> None:
    """Add a rounded-rectangle placeholder when no Market Intelligence chart data."""
    chart_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    chart_fill = chart_shape.fill
    chart_fill.solid()
    chart_fill.fore_color.rgb = RGBColor(15, 23, 42)
    chart_line = chart_shape.line
    chart_line.color.rgb = RGBColor(0, 255, 229)
    chart_line.width = Pt(1.5)
    tf = chart_shape.text_frame
    tf.clear()
    p_title = tf.paragraphs[0]
    p_title.text = "TAM / SAM / SOM"
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(224, 231, 255)
    p_title.alignment = PP_ALIGN.CENTER
    p_sub = tf.add_paragraph()
    p_sub.text = "Run Market Intelligence to see chart"
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = RGBColor(148, 163, 184)
    p_sub.alignment = PP_ALIGN.CENTER


# Slide index to narrative chapter id (hook=0, problem=1, solution=2, business=3, market=4)
_CHAPTER_IDS = ["hook", "problem", "solution", "business", "market"]


def _hex_to_rgb(hex_str: str) -> tuple[int, int, int] | None:
    """Convert #RRGGBB or RRGGBB to (r, g, b) 0-255. Returns None if invalid."""
    if not hex_str or not isinstance(hex_str, str):
        return None
    s = hex_str.strip().lstrip("#")
    if len(s) != 6:
        return None
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return None


def build_pitch_deck(
    slides_data: List[dict],
    output_path: Path,
    market_intel: Any = None,
    image_placements: Any = None,
    style_overrides: Any = None,
) -> None:
    """
    Create a PowerPoint file from slides_data.
    Each item: {"title": str, "bullets": [str, ...]}.
    If market_intel has market_share_data, the market slide gets a real bar chart.
    If image_placements is a list of {"path": str, "slide_id": str}, those images are embedded on the given slides.
    If style_overrides is a dict with title_font_size_pt, bullet_font_size_pt, title_color_hex, bullet_color_hex,
    title_slide_bg_hex, content_slide_bg_hex, those are applied to the deck.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    placements_by_slide = {}
    if isinstance(image_placements, list):
        for p in image_placements:
            if isinstance(p, dict) and p.get("path") and p.get("slide_id"):
                sid = p["slide_id"]
                placements_by_slide.setdefault(sid, []).append(p["path"])

    so = style_overrides if isinstance(style_overrides, dict) else {}
    # Resolve style with defaults (hex -> RGB tuples for RGBColor)
    title_slide_bg = _hex_to_rgb(str(so.get("title_slide_bg_hex", ""))) or (2, 6, 23)
    content_slide_bg = _hex_to_rgb(str(so.get("content_slide_bg_hex", ""))) or (15, 23, 42)
    title_color = _hex_to_rgb(str(so.get("title_color_hex", ""))) or (224, 231, 255)
    bullet_color = _hex_to_rgb(str(so.get("bullet_color_hex", ""))) or (226, 232, 240)
    title_slide_font_pt = int(so["title_font_size_pt"]) if isinstance(so.get("title_font_size_pt"), (int, float)) else 44
    content_title_font_pt = int(so["title_font_size_pt"]) if isinstance(so.get("title_font_size_pt"), (int, float)) else 32
    bullet_font_pt = int(so["bullet_font_size_pt"]) if isinstance(so.get("bullet_font_size_pt"), (int, float)) else 14
    subtitle_color = _hex_to_rgb(str(so.get("bullet_color_hex", ""))) or (148, 163, 184)

    for i, slide_spec in enumerate(slides_data):
        title = (slide_spec.get("title") or "Slide").strip()
        bullets = slide_spec.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = [str(bullets)] if bullets else []

        # Title slide layout for first slide, title and content for rest
        if i == 0:
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            try:
                slide.follow_master_background = False
            except Exception:
                pass
            bg_fill = slide.background.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = RGBColor(*title_slide_bg)

            left = Inches(0.5)
            top = Inches(2.5)
            width = prs.slide_width - Inches(1)
            height = Inches(1.2)
            tx = slide.shapes.add_textbox(left, top, width, height)
            tf = tx.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(title_slide_font_pt)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_color)
            p.alignment = 1  # center
            if bullets:
                sub = slide.shapes.add_textbox(left, Inches(3.8), width, Inches(1.5))
                sub.text_frame.word_wrap = True
                sub_p = sub.text_frame.paragraphs[0]
                sub_p.text = bullets[0] if isinstance(bullets[0], str) else str(bullets[0])
                sub_p.font.size = Pt(min(bullet_font_pt + 4, 24))
                sub_p.font.color.rgb = RGBColor(*subtitle_color)
                sub_p.alignment = 1
            # Title slide can also have an embedded image (e.g. logo/screenshot)
            for img_path in placements_by_slide.get("hook", [])[:1]:
                try:
                    resolved = Path(img_path).resolve()
                    if not resolved.is_file():
                        continue
                    pic_left = prs.slide_width - Inches(4.0)
                    pic_top = Inches(4.2)
                    slide.shapes.add_picture(str(resolved), pic_left, pic_top, width=Inches(3.5))
                except Exception as e:
                    import logging
                    logging.getLogger(__name__).warning("Could not embed image on title slide: %s", e)
        else:
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            try:
                slide.follow_master_background = False
            except Exception:
                pass
            bg_fill = slide.background.fill
            bg_fill.solid()
            bg_fill.fore_color.rgb = RGBColor(*content_slide_bg)

            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
            title_p = title_box.text_frame.paragraphs[0]
            title_p.text = title
            title_p.font.size = Pt(content_title_font_pt)
            title_p.font.bold = True
            title_p.font.color.rgb = RGBColor(*title_color)

            lower_title = title.lower()
            all_bullets_text = " ".join(str(b) for b in bullets).lower()
            is_market_slide = any(key in lower_title for key in ("market", "tam", "sam", "som")) or any(
                key in all_bullets_text for key in ("market", "tam", "sam", "som", "cagr")
            )
            market_share_data = None
            if is_market_slide and market_intel and getattr(market_intel, "market_share_data", None):
                raw = getattr(market_intel, "market_share_data", None)
                if isinstance(raw, list) and len(raw) > 0:
                    market_share_data = raw

            # Bullets: use half width when slide has chart or embedded image (so content doesn't overlap)
            slide_id = _CHAPTER_IDS[i] if i < len(_CHAPTER_IDS) else None
            has_image = slide_id and slide_id in placements_by_slide
            body_width = Inches(7.0) if ((is_market_slide and market_share_data) or has_image) else Inches(12)
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), body_width, Inches(5.5))
            body_box.text_frame.word_wrap = True
            for j, bullet in enumerate(bullets[:8]):  # max 8 bullets per slide
                text = (bullet if isinstance(bullet, str) else str(bullet)).strip()
                if not text:
                    continue
                if j == 0:
                    p = body_box.text_frame.paragraphs[0]
                else:
                    p = body_box.text_frame.add_paragraph()
                p.text = f"• {text}"
                p.font.size = Pt(bullet_font_pt)
                p.font.color.rgb = RGBColor(*bullet_color)
                p.space_after = Pt(6)

            # Market slide: real chart from Market Intelligence or placeholder
            if is_market_slide:
                chart_left = prs.slide_width - Inches(5.0)
                chart_top = Inches(1.4)
                chart_width = Inches(4.8)
                chart_height = Inches(4.2)
                if market_share_data and _CHARTS_AVAILABLE and CategoryChartData is not None and XL_CHART_TYPE is not None:
                    try:
                        chart_data = CategoryChartData()
                        names = []
                        shares = []
                        for item in market_share_data[:10]:
                            if isinstance(item, dict):
                                n = item.get("name") or item.get("segment") or "?"
                                s = item.get("share")
                            else:
                                n = getattr(item, "name", None) or getattr(item, "segment", None) or "?"
                                s = getattr(item, "share", None)
                            if n is not None and s is not None:
                                names.append(str(n))
                                shares.append(float(s))
                        if names and shares:
                            chart_data.categories = names
                            chart_data.add_series("Share %", tuple(shares))
                            graphic_frame = slide.shapes.add_chart(
                                XL_CHART_TYPE.COLUMN_CLUSTERED,
                                chart_left, chart_top, chart_width, chart_height,
                                chart_data,
                            )
                            chart = graphic_frame.chart
                            chart.has_legend = False
                            value_axis = chart.value_axis
                            if value_axis.tick_labels:
                                value_axis.tick_labels.number_format = "0\"%\""
                    except Exception:
                        _add_market_chart_placeholder(slide, chart_left, chart_top, chart_width, chart_height)
                else:
                    _add_market_chart_placeholder(slide, chart_left, chart_top, chart_width, chart_height)

            # Embed image on this slide if requested (right side, proper alignment; one per slide)
            if i < len(_CHAPTER_IDS):
                slide_id = _CHAPTER_IDS[i]
                for img_path in placements_by_slide.get(slide_id, [])[:1]:
                    try:
                        resolved = Path(img_path).resolve()
                        if not resolved.is_file():
                            continue
                        # Right side: leave margin; image width 4.8", height auto to preserve aspect ratio
                        pic_left = prs.slide_width - Inches(5.2)
                        pic_top = Inches(1.35)
                        pic_width = Inches(4.6)
                        slide.shapes.add_picture(
                            str(resolved),
                            pic_left,
                            pic_top,
                            width=pic_width,
                        )
                    except Exception as e:
                        import logging
                        logging.getLogger(__name__).warning("Could not embed image on slide: %s", e)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def parse_slides_from_json(text: str) -> SlideResult:
    """Extract slides array and optional golden_thread from Nova response (may be wrapped in ```json ... ```)."""
    raw = text.strip()
    if "```json" in raw:
        start = raw.index("```json") + 7
        end = raw.find("```", start)
        raw = (raw[start:end] if end >= 0 else raw[start:]).strip()
    elif "```" in raw:
        start = raw.rfind("```")
        raw = raw[start:].replace("```", "").strip()
    if not raw.startswith("{"):
        return (None, None)
    try:
        data = json.loads(raw)
        if not isinstance(data, dict):
            return (None, None)
        slides = data.get("slides")
        golden = data.get("golden_thread")
        if isinstance(golden, str):
            golden = golden.strip() or None
        if isinstance(slides, list) and len(slides) > 0:
            return (slides, golden)
        return (None, golden)
    except Exception:
        return (None, None)
